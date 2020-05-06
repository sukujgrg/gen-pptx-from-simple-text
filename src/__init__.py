#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import random
import re
from pathlib import Path

import click
import pptx

def create_presentation(lyrics_file: Path,
                        pptx_template: Path,
                        master_slide_idx: int,
                        slide_layout_idx: int,
                        font_size: int,
                        dst_dir: Path,
                        font_name: str,
                        slide_txt_alignment: str = "left"
                        ) -> Path:
    prs = pptx.Presentation(pptx_template)
    # setting text box size and position
    slide_height = pptx.util.Length(prs.slide_height)
    slide_width = pptx.util.Length(prs.slide_width)
    # Emu is English metric unit
    tb_pos = {
        "left": pptx.util.Emu(400000),
        "top": pptx.util.Emu(400000),
        "width": pptx.util.Emu(slide_width - (400000 * 2)),
        "height": pptx.util.Emu(slide_height - (400000 * 2))
    }
    slide_layout = prs.slide_masters[master_slide_idx].slide_layouts[slide_layout_idx]
    with lyrics_file.open() as f:
        emptyline_re = r"(?:\r?\n){2,}"
        lyrics = re.split(emptyline_re, f.read())
        for content in lyrics:
            slide = prs.slides.add_slide(slide_layout)
            txbox = slide.shapes.add_textbox(**tb_pos)
            tf = txbox.text_frame
            # this is to keep text frame in the middle of the screen from top to bottom of the screen
            tf.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
            # this is supposed to work as per the documentation but it's not working
            # 09/12/2018
            # tf.fit_text(font_family=font_name, max_size=font_size)
            p = tf.add_paragraph()
            if slide_txt_alignment == "left":
                p.alignment = pptx.enum.text.PP_ALIGN.LEFT
            elif slide_txt_alignment == "middle":
                p.alignment = pptx.enum.text.PP_ALIGN.CENTER
            else:
                p.alignment = pptx.enum.text.PP_ALIGN.RIGHT
            p.font.name = font_name
            p.font.size = pptx.util.Pt(font_size)
            p.text = content
    dst_dir.mkdir(parents=True, exist_ok=True)
    dst_file = lyrics_file.with_suffix(".pptx")
    dst_path = Path(dst_dir, dst_file.name)
    if dst_path.exists():
        dst_path.unlink()
    prs.save(str(dst_path))
    return dst_path


def pick_master_slide(master_slides_path: Path) -> Path:
    """Returns path to a master slide. If argument is a directory,
    returns a randomly chosen file in the directory and if argument
    is a file, returns the argument itself
    """
    if master_slides_path.is_file():
        return master_slides_path
    master_slides = [
        f.name for f in master_slides_path.iterdir()
        if f.suffix == '.pptx'
        if not f.name.startswith("archived")
    ]
    if master_slides:
        return Path(master_slides_path, random.choice(master_slides))
    raise ValueError(f"Unable to find any valid pptx template file in {master_slides_path}")


@click.command()
@click.argument("lyrics-files", nargs=-1, type=click.Path(exists=True, dir_okay=False), required=True)
@click.option("--pptx-template-path", "-pt", type=click.Path(exists=True), required=True)
@click.option("--master-slide-idx", "-ms", default=0, type=int, show_default=True)
@click.option("--slide-layout-idx", "-sl", default=6, type=int, show_default=True)
@click.option("--font-size", "-fs", default=32, type=int, show_default=True)
@click.option("--font-name", "-fn", default="Calibri", type=str, show_default=True)
@click.option("--dst-dir", "-ns", default="./generated-pptx", show_default=True)
@click.option("--slide-txt-alignment", "-ta", default="left",
              type=click.Choice(["left", "middle", "right"]),
              show_default=True)
def cli(lyrics_files, pptx_template_path, font_size, master_slide_idx, slide_layout_idx, dst_dir, font_name,
        slide_txt_alignment):
    """A powerpoint generator.
    
    https://github.com/sukujgrg/gen-pptx-from-simple-text
    """
    dst_dir = Path(dst_dir)
    pptx_template_path = Path(pptx_template_path)
    pptx_template = pick_master_slide(pptx_template_path)
    for f in lyrics_files:
        try:
            r = create_presentation(
                    Path(f),
                    pptx_template,
                    master_slide_idx,
                    slide_layout_idx,
                    font_size,
                    dst_dir,
                    font_name,
                    slide_txt_alignment
                )
            msg = f"PPTX: {r}"
            click.echo(click.style(msg, fg="green"))
        except Exception:
            raise


if __name__ == '__main__':
    cli()
